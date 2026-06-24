#!/usr/bin/env python3
"""Branded PowerPoint deck (also imports cleanly into Google Slides).

Mirrors the beamer deck in the fastverk palette. Text uses Space Grotesk (renders
if installed on the viewer; falls back gracefully). Imagery is the same icons the
rest of the pipeline produces.

CLI: pptx_gen.py <out.pptx> <wordmark_dark.png> <icon_dark.png> <icon_light.png>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG, FG = RGBColor(0x15, 0x16, 0x1A), RGBColor(0xEC, 0xE7, 0xDA)
AMBER, SLATE = RGBColor(0xE0, 0xA3, 0x3E), RGBColor(0x4A, 0x56, 0x5A)
FONT = "Space Grotesk"
SW, SH = Inches(13.333), Inches(7.5)

def _bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _text(slide, text, left, top, width, height, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb

def _center_pic(slide, path, top, height):
    pic = slide.shapes.add_picture(path, 0, top, height=height)
    pic.left = int((SW - pic.width) / 2)
    return pic

def build(out, wordmark, icon_dark, icon_light):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # title
    s = prs.slides.add_slide(blank); _bg(s)
    _center_pic(s, wordmark, Inches(2.4), Inches(1.7))
    _text(s, "Brand deck", 0, Inches(4.5), SW, Inches(0.6), 22, SLATE, align=PP_ALIGN.CENTER)
    _text(s, "generated from one parametric source", 0, Inches(5.1), SW, Inches(0.5), 14, SLATE, align=PP_ALIGN.CENTER)

    # the mark
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, "The mark", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, bold=True)
    _text(s, "Plays off F, V, the down-triangle and the arms\nA golden-ratio construction on the unit circle\nOne source for SVG, icons, brandbook, decks",
          Inches(0.9), Inches(1.8), Inches(7), Inches(4), 20, FG)
    s.shapes.add_picture(icon_dark, Inches(9.2), Inches(2.0), height=Inches(3.2))

    # light & dark
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, "Light & dark", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, bold=True)
    s.shapes.add_picture(icon_dark, Inches(3.0), Inches(2.2), height=Inches(3.0))
    s.shapes.add_picture(icon_light, Inches(7.8), Inches(2.2), height=Inches(3.0))
    _text(s, "dark", Inches(3.0), Inches(5.4), Inches(2.6), Inches(0.5), 16, SLATE, align=PP_ALIGN.CENTER)
    _text(s, "light", Inches(7.8), Inches(5.4), Inches(2.6), Inches(0.5), 16, SLATE, align=PP_ALIGN.CENTER)

    # color
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, "Color", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, bold=True)
    for i, (col, name) in enumerate([(BG, "15161A"), (FG, "ECE7DA"), (AMBER, "E0A33E"), (SLATE, "4A565A")]):
        x = Inches(0.9 + i * 3.0)
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(2.4), Inches(1.7))
        sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.color.rgb = SLATE
        _text(s, name, x, Inches(4.0), Inches(2.4), Inches(0.5), 14, FG)

    # closing
    s = prs.slides.add_slide(blank); _bg(s)
    _center_pic(s, icon_dark, Inches(1.8), Inches(2.6))
    _text(s, "fastverk", 0, Inches(5.0), SW, Inches(0.8), 28, FG, bold=True, align=PP_ALIGN.CENTER)

    prs.save(out)
    print("wrote", out)

if __name__ == "__main__":
    build(sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4])
